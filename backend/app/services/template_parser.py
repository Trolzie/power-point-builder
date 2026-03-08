from pathlib import Path

from pptx import Presentation

from app.models.template import (
    LayoutInfo,
    MasterInfo,
    PlaceholderInfo,
    TemplateManifest,
)

EMU_PER_PT = 12700
AVG_CHARS_PER_WORD = 5.5

SKIP_PLACEHOLDER_TYPES = {"DATE", "FOOTER", "SLIDE_NUMBER", "HEADER", "CHART", "TABLE", "VERTICAL_OBJECT", "VERTICAL_BODY"}


def _extract_placeholder_formatting(ph) -> dict:
    """Extract default font formatting from a placeholder's first paragraph/run."""
    result = {}
    try:
        tf = ph.text_frame
        if not tf.paragraphs:
            return result

        first_para = tf.paragraphs[0]

        if first_para.alignment is not None:
            result["default_alignment"] = str(first_para.alignment).split(".")[-1]

        if first_para.runs:
            font = first_para.runs[0].font
            if font.name:
                result["default_font_name"] = font.name
            if font.size is not None:
                result["default_font_size_pt"] = font.size.pt
            if font.bold is not None:
                result["default_font_bold"] = font.bold
            try:
                if font.color and font.color.rgb:
                    result["default_font_color"] = str(font.color.rgb)
            except (AttributeError, TypeError):
                pass
    except Exception:
        pass
    return result


def _estimate_capacity(width_emu, height_emu, font_size_pt, line_spacing=1.15):
    """Estimate how many lines/words fit in a placeholder."""
    if not all([width_emu, height_emu, font_size_pt]):
        return None, None

    font_size_emu = font_size_pt * EMU_PER_PT
    line_height = font_size_emu * line_spacing

    max_lines = max(1, int(height_emu / line_height))
    avg_char_width = font_size_emu * 0.55
    chars_per_line = int(width_emu / avg_char_width) if avg_char_width > 0 else 50
    max_words = max(1, int((max_lines * chars_per_line) / AVG_CHARS_PER_WORD))

    return max_lines, max_words


def _placeholder_type_str(ph_type) -> str:
    """Convert placeholder type enum to string."""
    if ph_type is None:
        return "UNKNOWN"
    # python-pptx enums have a .name property
    if hasattr(ph_type, "name") and ph_type.name:
        return ph_type.name
    return str(ph_type)


def _extract_theme_colors(prs: Presentation) -> dict[str, str]:
    """Extract theme colors from the first slide master's theme XML."""
    colors: dict[str, str] = {}
    try:
        theme_element = prs.slide_masters[0].element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        clr_schemes = theme_element.findall(".//a:clrScheme", ns)
        if not clr_schemes:
            return colors

        clr_scheme = clr_schemes[0]
        color_names = [
            "dk1", "dk2", "lt1", "lt2",
            "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
            "hlink", "folHlink",
        ]
        for name in color_names:
            elem = clr_scheme.find(f"a:{name}", ns)
            if elem is not None:
                # Color value can be in srgbClr or sysClr child
                srgb = elem.find("a:srgbClr", ns)
                if srgb is not None:
                    colors[name] = srgb.get("val", "")
                else:
                    sys_clr = elem.find("a:sysClr", ns)
                    if sys_clr is not None:
                        colors[name] = sys_clr.get("lastClr", sys_clr.get("val", ""))
    except Exception:
        pass
    return colors


def parse_template(file_path: str, template_id: str | None = None) -> TemplateManifest:
    """Parse a .pptx template and extract its structure as a manifest."""
    path = Path(file_path)
    prs = Presentation(str(path))

    if template_id is None:
        template_id = path.stem

    masters: list[MasterInfo] = []
    for master_idx, master in enumerate(prs.slide_masters):
        layouts: list[LayoutInfo] = []
        for layout_idx, layout in enumerate(master.slide_layouts):
            placeholders: list[PlaceholderInfo] = []
            for ph in layout.placeholders:
                fmt = _extract_placeholder_formatting(ph)
                font_size = fmt.get("default_font_size_pt")
                max_lines, max_words = _estimate_capacity(
                    ph.width, ph.height, font_size or 18.0,
                )
                placeholders.append(
                    PlaceholderInfo(
                        idx=ph.placeholder_format.idx,
                        name=ph.name,
                        type=_placeholder_type_str(ph.placeholder_format.type),
                        left=ph.left,
                        top=ph.top,
                        width=ph.width,
                        height=ph.height,
                        estimated_max_lines=max_lines,
                        estimated_max_words=max_words,
                        **fmt,
                    )
                )
            content_phs = [
                ph for ph in placeholders
                if ph.type not in SKIP_PLACEHOLDER_TYPES
            ]
            content_count = len(content_phs)
            recommended = 1 <= content_count <= 8

            # Build preview description like "1 title + 2 body + 1 image"
            type_counts: dict[str, int] = {}
            for ph in content_phs:
                label = ph.type.lower()
                type_counts[label] = type_counts.get(label, 0) + 1
            preview = " + ".join(
                f"{count} {label}" for label, count in type_counts.items()
            ) if type_counts else None

            layouts.append(
                LayoutInfo(
                    index=layout_idx,
                    name=layout.name,
                    placeholders=placeholders,
                    content_placeholder_count=content_count,
                    recommended=recommended,
                    preview_description=preview,
                )
            )
        masters.append(
            MasterInfo(
                index=master_idx,
                name=f"Master {master_idx}",
                layouts=layouts,
            )
        )

    theme_colors = _extract_theme_colors(prs)

    return TemplateManifest(
        template_id=template_id,
        filename=path.name,
        slide_width_emu=prs.slide_width,
        slide_height_emu=prs.slide_height,
        theme_colors=theme_colors,
        masters=masters,
    )
