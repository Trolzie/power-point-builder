import copy
import logging
from pathlib import Path

from pptx import Presentation

from app.models.presentation import PlaceholderContent, PresentationContent

logger = logging.getLogger(__name__)

_NSMAP = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def _qn(tag: str) -> str:
    """Resolve a namespace-prefixed tag like 'a:rPr' to its Clark notation."""
    prefix, local = tag.split(":")
    return f"{{{_NSMAP[prefix]}}}{local}"


def _fill_placeholder(placeholder, content: PlaceholderContent) -> None:
    """Fill a placeholder with content while preserving formatting."""
    if content.type == "image" and content.image_path:
        if hasattr(placeholder, 'insert_picture'):
            placeholder.insert_picture(content.image_path)
        else:
            logger.warning(
                "Placeholder '%s' is not a picture placeholder, skipping image",
                placeholder.name,
            )
        return

    if content.type != "text" or not content.paragraphs:
        return

    tf = placeholder.text_frame
    template_para = tf.paragraphs[0]

    # Capture template formatting from first paragraph/run for reuse on added paragraphs
    template_pPr = template_para._p.find(_qn("a:pPr"))
    template_rPr = None
    if template_para.runs:
        template_rPr = template_para.runs[0]._r.find(_qn("a:rPr"))

    for i, para_content in enumerate(content.paragraphs):
        if i == 0:
            para = tf.paragraphs[0]
        else:
            para = tf.add_paragraph()
            # Copy paragraph-level properties (font defaults, spacing) from template
            if template_pPr is not None:
                new_pPr = copy.deepcopy(template_pPr)
                para._p.insert(0, new_pPr)

        para.level = para_content.level

        # Reuse existing run if available, otherwise add new one
        if para.runs:
            run = para.runs[0]
            run.text = para_content.text
            # Remove extra runs from paragraph XML
            for extra_run in para.runs[1:]:
                extra_run._r.getparent().remove(extra_run._r)
        else:
            run = para.add_run()
            run.text = para_content.text
            # Copy run-level properties (font, size, color) from template
            if template_rPr is not None:
                new_rPr = copy.deepcopy(template_rPr)
                run._r.insert(0, new_rPr)

        if para_content.bold is not None:
            run.font.bold = para_content.bold
        if para_content.italic is not None:
            run.font.italic = para_content.italic


def assemble_presentation(
    template_path: str, content: PresentationContent, output_path: str
) -> str:
    """Assemble a .pptx from a template and generated content."""
    prs = Presentation(template_path)

    # Remove any existing slides from the template (keep layouts/masters)
    for slide in list(prs.slides):
        for rel_key, rel in prs.part.rels.items():
            if rel.target_part is slide.part:
                prs.part.drop_rel(rel_key)
                break
    sldId_list = prs.slides._sldIdLst
    for sldId in list(sldId_list):
        sldId_list.remove(sldId)

    for slide_content in content.slides:
        layout = prs.slide_layouts[slide_content.layout_index]
        slide = prs.slides.add_slide(layout)

        for ph_key, ph_content in slide_content.placeholders.items():
            ph_idx = int(ph_key)
            # Find the placeholder on the slide by idx
            target_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == ph_idx:
                    target_ph = ph
                    break

            if target_ph is not None:
                _fill_placeholder(target_ph, ph_content)
            else:
                logger.warning(
                    "Placeholder idx=%s not found on slide layout '%s'",
                    ph_idx,
                    slide_content.layout_name,
                )

        if slide_content.speaker_notes:
            slide.notes_slide.notes_text_frame.text = slide_content.speaker_notes

    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path
